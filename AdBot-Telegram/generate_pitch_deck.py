"""
generate_pitch_deck.py — Generates AdBot Investor Pitch Deck (.pptx)
Run: python generate_pitch_deck.py
Output: docs/AdBot_Investor_Pitch_Deck.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt
import pptx.util as util
from pptx.enum.dml import MSO_THEME_COLOR
import os

# ── Brand Colors ──────────────────────────────────────────────────────────────
DEEP_NAVY    = RGBColor(0x0A, 0x16, 0x28)   # #0A1628 — primary background
ELECTRIC     = RGBColor(0x4F, 0x8E, 0xF7)   # #4F8EF7 — accent blue
GOLD         = RGBColor(0xF5, 0xA6, 0x23)   # #F5A623 — highlight/CTA
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)   # #FFFFFF
LIGHT_GREY   = RGBColor(0xB0, 0xBC, 0xCC)   # #B0BCCC — secondary text
MID_NAVY     = RGBColor(0x14, 0x2A, 0x4A)   # #142A4A — card background
GREEN        = RGBColor(0x2E, 0xCC, 0x71)   # #2ECC71 — positive metrics
RED_SOFT     = RGBColor(0xE7, 0x4C, 0x3C)   # for competitor weaknesses

W = Inches(13.33)   # Slide width  (16:9 widescreen)
H = Inches(7.5)     # Slide height


# ── Helper functions ──────────────────────────────────────────────────────────

def add_slide(prs, layout_idx=6):
    """Add a blank slide."""
    layout = prs.slide_layouts[layout_idx]
    return prs.slides.add_slide(layout)


def bg(slide, color=DEEP_NAVY):
    """Fill slide background with a solid color."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def box(slide, x, y, w, h, fill_color=None, border_color=None, border_width=Pt(0)):
    """Add a rectangle shape."""
    shape = slide.shapes.add_shape(1, x, y, w, h)  # 1 = MSO_SHAPE_TYPE.RECTANGLE
    shape.line.width = border_width
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_width
    else:
        shape.line.fill.background()
    return shape


def txt(slide, text, x, y, w, h,
        size=Pt(18), bold=False, color=WHITE, align=PP_ALIGN.LEFT,
        italic=False, wrap=True):
    """Add a text box."""
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = size
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox


def label_value(slide, label, value, x, y, w=Inches(2.8), label_color=LIGHT_GREY, value_color=WHITE, value_size=Pt(28)):
    """Add a metric card."""
    box(slide, x, y, w, Inches(1.3), fill_color=MID_NAVY)
    txt(slide, value, x + Inches(0.15), y + Inches(0.05), w - Inches(0.3), Inches(0.7),
        size=value_size, bold=True, color=value_color, align=PP_ALIGN.CENTER)
    txt(slide, label, x + Inches(0.15), y + Inches(0.75), w - Inches(0.3), Inches(0.45),
        size=Pt(11), color=label_color, align=PP_ALIGN.CENTER)


def accent_line(slide, x, y, w, color=ELECTRIC):
    """Thin colored horizontal line."""
    shape = slide.shapes.add_shape(1, x, y, w, Pt(3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def slide_number(slide, n, total=15):
    txt(slide, f"{n} / {total}", Inches(12.5), Inches(7.1), Inches(0.7), Inches(0.3),
        size=Pt(9), color=LIGHT_GREY, align=PP_ALIGN.RIGHT)


def bullet_list(slide, items, x, y, w, h, size=Pt(14), dot_color=ELECTRIC, text_color=WHITE):
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(4)
        run = p.add_run()
        run.text = f"  ▸  {item}"
        run.font.size = size
        run.font.color.rgb = text_color


# ── Slide Builders ────────────────────────────────────────────────────────────

def slide_cover(prs):
    s = add_slide(prs)
    bg(s, DEEP_NAVY)

    # Gradient accent bar top
    box(s, 0, 0, W, Inches(0.08), fill_color=ELECTRIC)

    # Gold accent
    box(s, 0, Inches(0.08), Inches(0.6), H - Inches(0.08), fill_color=GOLD)

    # Main headline
    txt(s, "AdBot", Inches(1), Inches(1.2), Inches(8), Inches(1.5),
        size=Pt(80), bold=True, color=WHITE)
    txt(s, "The AI Marketing Agency", Inches(1), Inches(2.7), Inches(9), Inches(0.8),
        size=Pt(32), bold=False, color=ELECTRIC)
    txt(s, "in Your Pocket.", Inches(1), Inches(3.3), Inches(9), Inches(0.8),
        size=Pt(32), bold=True, color=GOLD)

    accent_line(s, Inches(1), Inches(4.3), Inches(5), ELECTRIC)

    txt(s, "Professional AI-generated ads in any language\ndelivered via WhatsApp & Telegram — for free.",
        Inches(1), Inches(4.5), Inches(9), Inches(1.2),
        size=Pt(17), color=LIGHT_GREY)

    txt(s, "SEED ROUND   |   2026", Inches(1), Inches(6.6), Inches(5), Inches(0.5),
        size=Pt(13), color=GOLD, bold=True)


def slide_problem(prs):
    s = add_slide(prs)
    bg(s, DEEP_NAVY)
    box(s, 0, 0, W, Inches(0.08), fill_color=ELECTRIC)
    slide_number(s, 2)

    txt(s, "THE PROBLEM", Inches(0.6), Inches(0.3), Inches(5), Inches(0.5),
        size=Pt(11), color=ELECTRIC, bold=True)
    txt(s, "400 Million Small Businesses.\nZero Affordable Ad Tools.",
        Inches(0.6), Inches(0.7), Inches(8), Inches(1.6),
        size=Pt(36), bold=True, color=WHITE)

    accent_line(s, Inches(0.6), Inches(2.4), Inches(4), GOLD)

    problems = [
        ("💸", "Ad agencies cost $500-5000/month", "Unaffordable for 95% of small businesses"),
        ("🌐", "AdCreative.ai & Canva require a website", "Shop owners in Mumbai & Lagos use WhatsApp, not browsers"),
        ("🗣️", "All tools are English-first", "1.9B South Asians, 500M Arabic speakers are ignored"),
        ("📸", "Text-to-image AI generates wrong products", "The shoe looks nothing like your actual shoe"),
    ]

    for i, (icon, headline, sub) in enumerate(problems):
        col = Inches(0.5) if i % 2 == 0 else Inches(7)
        row = Inches(2.7) if i < 2 else Inches(4.7)
        box(s, col, row, Inches(5.8), Inches(1.7), fill_color=MID_NAVY)
        txt(s, icon + "  " + headline, col + Inches(0.2), row + Inches(0.15),
            Inches(5.4), Inches(0.5), size=Pt(15), bold=True, color=WHITE)
        txt(s, sub, col + Inches(0.2), row + Inches(0.65),
            Inches(5.4), Inches(0.8), size=Pt(12), color=LIGHT_GREY)


def slide_solution(prs):
    s = add_slide(prs)
    bg(s, DEEP_NAVY)
    box(s, 0, 0, W, Inches(0.08), fill_color=GOLD)
    slide_number(s, 3)

    txt(s, "THE SOLUTION", Inches(0.6), Inches(0.3), Inches(5), Inches(0.5),
        size=Pt(11), color=GOLD, bold=True)
    txt(s, "AdBot — Photo → Professional Ad\nin 90 Seconds. In Your Language.",
        Inches(0.6), Inches(0.7), Inches(12), Inches(1.5),
        size=Pt(32), bold=True, color=WHITE)

    accent_line(s, Inches(0.6), Inches(2.3), Inches(6), ELECTRIC)

    # Pipeline steps
    steps = [
        ("📸", "Upload\nProduct Photo", ELECTRIC),
        ("🧠", "Gemini AI\nAnalyzes & Writes", GOLD),
        ("✂️", "Auto Remove\nBackground", ELECTRIC),
        ("🎨", "Compose\nProfessional Ad", GOLD),
        ("📲", "Receive in\nWhatsApp/Telegram", GREEN),
    ]

    for i, (icon, label, color) in enumerate(steps):
        x = Inches(0.4) + i * Inches(2.58)
        box(s, x, Inches(2.7), Inches(2.3), Inches(2.0), fill_color=MID_NAVY)
        txt(s, icon, x + Inches(0.85), Inches(2.8), Inches(0.6), Inches(0.6), size=Pt(26))
        txt(s, label, x + Inches(0.1), Inches(3.4), Inches(2.1), Inches(0.9),
            size=Pt(13), bold=True, color=color, align=PP_ALIGN.CENTER)
        if i < 4:
            txt(s, "→", x + Inches(2.3), Inches(3.35), Inches(0.3), Inches(0.4),
                size=Pt(20), color=LIGHT_GREY, align=PP_ALIGN.CENTER)

    features = [
        "✅  Real product photo (not hallucinated)",
        "✅  50+ languages including Hindi, Arabic, Swahili",
        "✅  Headline + CTA text overlaid on the image (like AdCreative.ai)",
        "✅  Works entirely in WhatsApp & Telegram — no app, no website",
        "✅  $0 operating cost — free for users, monetised via PRO tier",
    ]
    for i, f in enumerate(features):
        txt(s, f, Inches(0.6), Inches(4.9) + i * Inches(0.44),
            Inches(12), Inches(0.4), size=Pt(13), color=WHITE)


def slide_market(prs):
    s = add_slide(prs)
    bg(s, DEEP_NAVY)
    box(s, 0, 0, W, Inches(0.08), fill_color=ELECTRIC)
    slide_number(s, 4)

    txt(s, "MARKET OPPORTUNITY", Inches(0.6), Inches(0.3), Inches(6), Inches(0.5),
        size=Pt(11), color=ELECTRIC, bold=True)
    txt(s, "A $150B Market That's Completely Underserved",
        Inches(0.6), Inches(0.7), Inches(11), Inches(0.9),
        size=Pt(30), bold=True, color=WHITE)

    # TAM / SAM / SOM
    markets = [
        ("TAM", "$740B", "Global Digital\nAd Market (2024)", WHITE),
        ("SAM", "$150B", "SMB Digital\nAd Spend", ELECTRIC),
        ("SOM", "$2.4B", "WhatsApp/Telegram\nSMB Tools (5yr)", GOLD),
    ]
    for i, (label, val, desc, color) in enumerate(markets):
        x = Inches(0.5) + i * Inches(4.25)
        box(s, x, Inches(1.9), Inches(3.8), Inches(2.4), fill_color=MID_NAVY)
        txt(s, label, x + Inches(0.2), Inches(2.0), Inches(3.4), Inches(0.5),
            size=Pt(13), bold=True, color=color, align=PP_ALIGN.CENTER)
        txt(s, val, x + Inches(0.1), Inches(2.45), Inches(3.6), Inches(0.9),
            size=Pt(42), bold=True, color=color, align=PP_ALIGN.CENTER)
        txt(s, desc, x + Inches(0.1), Inches(3.3), Inches(3.6), Inches(0.8),
            size=Pt(12), color=LIGHT_GREY, align=PP_ALIGN.CENTER)

    # Key stats
    stats = [
        ("400M+", "Small businesses\nglobally"),
        ("2.8B", "WhatsApp\nmonthly users"),
        ("85%", "SMBs can't afford\nad agencies"),
        ("50+", "Languages we\ncan serve"),
    ]
    for i, (val, desc) in enumerate(stats):
        x = Inches(0.5) + i * Inches(3.2)
        label_value(s, desc, val, x, Inches(4.6), w=Inches(2.9))


def slide_pipeline(prs):
    s = add_slide(prs)
    bg(s, DEEP_NAVY)
    box(s, 0, 0, W, Inches(0.08), fill_color=GOLD)
    slide_number(s, 5)

    txt(s, "HOW IT WORKS", Inches(0.6), Inches(0.3), Inches(5), Inches(0.5),
        size=Pt(11), color=GOLD, bold=True)
    txt(s, "The Vision DNA Pipeline",
        Inches(0.6), Inches(0.7), Inches(9), Inches(0.8),
        size=Pt(30), bold=True, color=WHITE)

    stages = [
        ("STAGE 1", "Product Isolation", "rembg removes background from user's photo.\nThe real product is extracted with 95%+ clean edges.\nCost: $0 — runs locally, no API.", ELECTRIC),
        ("STAGE 2", "AI Copy Generation", "Gemini 2.5 Flash analyzes the product + writes:\nHeadline, Body, CTA, Hashtags, Visual style brief.\nIn the user's native language. Cost: $0 free tier.", GOLD),
        ("STAGE 3", "Ad Composition", "Pillow compositor places product on pro background,\nrenders headline + CTA text as crisp vector-quality text,\napplies brand gradients + shadows. Cost: $0.", GREEN),
        ("STAGE 4", "Delivery", "Final ad image sent directly in Telegram or WhatsApp.\nUser also receives copy pack: Instagram, WhatsApp,\nGoogle Ad format, Print Poster. Cost: $0.", ELECTRIC),
    ]

    for i, (stage, title, desc, color) in enumerate(stages):
        x = Inches(0.4) if i % 2 == 0 else Inches(6.9)
        y = Inches(1.8) if i < 2 else Inches(4.5)
        box(s, x, y, Inches(6.2), Inches(2.4), fill_color=MID_NAVY)
        box(s, x, y, Inches(1.1), Inches(2.4), fill_color=color)
        txt(s, stage, x + Inches(0.05), y + Inches(0.9), Inches(1.0), Inches(0.6),
            size=Pt(10), bold=True, color=DEEP_NAVY, align=PP_ALIGN.CENTER)
        txt(s, title, x + Inches(1.2), y + Inches(0.15), Inches(4.8), Inches(0.5),
            size=Pt(16), bold=True, color=color)
        txt(s, desc, x + Inches(1.2), y + Inches(0.65), Inches(4.8), Inches(1.6),
            size=Pt(12), color=LIGHT_GREY)


def slide_competitive(prs):
    s = add_slide(prs)
    bg(s, DEEP_NAVY)
    box(s, 0, 0, W, Inches(0.08), fill_color=ELECTRIC)
    slide_number(s, 6)

    txt(s, "COMPETITIVE LANDSCAPE", Inches(0.6), Inches(0.3), Inches(6), Inches(0.5),
        size=Pt(11), color=ELECTRIC, bold=True)
    txt(s, "We Win on 5 Critical Dimensions",
        Inches(0.6), Inches(0.7), Inches(10), Inches(0.8),
        size=Pt(28), bold=True, color=WHITE)

    headers = ["Feature", "AdCreative.ai", "Canva", "DIY Agency", "AdBot ✦"]
    col_w = [Inches(3.2), Inches(2.1), Inches(2.1), Inches(2.1), Inches(2.1)]
    x_starts = [Inches(0.4), Inches(3.6), Inches(5.7), Inches(7.8), Inches(9.9)]

    # Header row
    for i, h in enumerate(headers):
        c = ELECTRIC if i == 4 else MID_NAVY
        box(s, x_starts[i], Inches(1.8), col_w[i] - Inches(0.05), Inches(0.5), fill_color=c)
        tc = DEEP_NAVY if i == 4 else WHITE
        txt(s, h, x_starts[i] + Inches(0.05), Inches(1.85), col_w[i] - Inches(0.1), Inches(0.4),
            size=Pt(12), bold=True, color=tc, align=PP_ALIGN.CENTER)

    rows = [
        ("Real product photo in ad",     "✅", "Manual", "✅", "✅"),
        ("Works in WhatsApp/Telegram",   "❌", "❌",     "❌", "✅"),
        ("50+ languages",                "❌", "Basic",  "❌", "✅"),
        ("Text overlay on image",        "✅", "Manual", "✅", "✅"),
        ("AI copy generation",           "✅", "❌",     "❌", "✅"),
        ("Price",                   "$29/mo",  "Free*", "$500+", "FREE"),
    ]

    for r, (feature, ac, cv, diy, us) in enumerate(rows):
        y = Inches(2.4) + r * Inches(0.72)
        alt = r % 2 == 1
        row_bg = RGBColor(0x12, 0x25, 0x40) if alt else MID_NAVY

        for i, val in enumerate([feature, ac, cv, diy, us]):
            c = RGBColor(0x1A, 0x3A, 0x60) if (i == 4 and alt) else (ELECTRIC if i == 4 else row_bg)
            box(s, x_starts[i], y, col_w[i] - Inches(0.05), Inches(0.68), fill_color=c)
            vc = GOLD if (i == 4 and r == 5) else (GREEN if (i == 4 and val == "✅") else
                 (RED_SOFT if val == "❌" else WHITE))
            txt(s, val, x_starts[i] + Inches(0.05), y + Inches(0.1),
                col_w[i] - Inches(0.1), Inches(0.48), size=Pt(12), color=vc,
                align=PP_ALIGN.CENTER if i > 0 else PP_ALIGN.LEFT, bold=(i == 4))


def slide_business_model(prs):
    s = add_slide(prs)
    bg(s, DEEP_NAVY)
    box(s, 0, 0, W, Inches(0.08), fill_color=GOLD)
    slide_number(s, 7)

    txt(s, "BUSINESS MODEL", Inches(0.6), Inches(0.3), Inches(5), Inches(0.5),
        size=Pt(11), color=GOLD, bold=True)
    txt(s, "4 Revenue Streams. All Running on Day 1.",
        Inches(0.6), Inches(0.7), Inches(11), Inches(0.8),
        size=Pt(30), bold=True, color=WHITE)

    streams = [
        ("⭐", "Telegram Stars PRO",
         ["Free tier: 3 ads/day", "PRO Weekly: 50 Stars (~$0.65)", "PRO Monthly: 150 Stars (~$1.95)",
          "PRO Quarterly: 500 Stars (~$6.50)", "Pay inside Telegram — no card needed"],
         ELECTRIC, "Recurring subscription revenue"),
        ("🎨", "Canva Affiliate",
         ["Every ad delivery includes Canva link", "Tagged with affiliate ID", "User signs up → 25-80% commission",
          "Zero effort — auto-appended", "~$30-50 lifetime value per conversion"],
         GOLD, "Passive affiliate income"),
        ("🏢", "WhatsApp Business API",
         ["Enterprise tier for larger SMBs", "Custom branding + bulk ad generation",
          "White-label for marketing agencies", "Volume pricing: $50-500/mo",
          "API access for developers"],
         GREEN, "B2B enterprise revenue"),
        ("📊", "Data & Insights (Year 2)",
         ["Anonymized ad performance trends", "Industry benchmarks for marketers",
          "Sold to marketing research firms", "Ad effectiveness scoring",
          "No user PII — fully ethical"],
         RGBColor(0xAA, 0x88, 0xFF), "Future data product"),
    ]

    for i, (icon, title, points, color, subtitle) in enumerate(streams):
        x = Inches(0.4) + (i % 2) * Inches(6.5)
        y = Inches(1.9) + (i // 2) * Inches(2.6)
        box(s, x, y, Inches(6.2), Inches(2.4), fill_color=MID_NAVY)
        box(s, x, y, Inches(0.12), Inches(2.4), fill_color=color)
        txt(s, icon + "  " + title, x + Inches(0.25), y + Inches(0.12),
            Inches(5.8), Inches(0.45), size=Pt(15), bold=True, color=color)
        txt(s, subtitle, x + Inches(0.25), y + Inches(0.55),
            Inches(5.8), Inches(0.3), size=Pt(10), color=LIGHT_GREY, italic=True)
        for j, pt in enumerate(points[:3]):
            txt(s, "▸  " + pt, x + Inches(0.25), y + Inches(0.9) + j * Inches(0.44),
                Inches(5.8), Inches(0.4), size=Pt(12), color=WHITE)


def slide_financials(prs):
    s = add_slide(prs)
    bg(s, DEEP_NAVY)
    box(s, 0, 0, W, Inches(0.08), fill_color=GREEN)
    slide_number(s, 8)

    txt(s, "FINANCIAL PROJECTIONS", Inches(0.6), Inches(0.3), Inches(7), Inches(0.5),
        size=Pt(11), color=GREEN, bold=True)
    txt(s, "Conservative 3-Year Revenue Model",
        Inches(0.6), Inches(0.7), Inches(11), Inches(0.8),
        size=Pt(30), bold=True, color=WHITE)

    # Assumptions note
    txt(s, "Assumptions: Free viral growth via WhatsApp sharing (K-factor 1.3). No paid acquisition in Year 1.",
        Inches(0.6), Inches(1.5), Inches(12), Inches(0.4),
        size=Pt(11), color=LIGHT_GREY, italic=True)

    years = [
        ("YEAR 1", "Launch & Grow",
         [("Total Users", "25,000"), ("PRO Users (5%)", "1,250"),
          ("Monthly Recurring Revenue", "$3,125"), ("Affiliate Income/mo", "$500"),
          ("Annual Revenue", "~$43,500")], ELECTRIC),
        ("YEAR 2", "Scale",
         [("Total Users", "200,000"), ("PRO Users (8%)", "16,000"),
          ("Monthly Recurring Revenue", "$40,000"), ("Affiliate Income/mo", "$6,000"),
          ("Annual Revenue", "~$552,000")], GOLD),
        ("YEAR 3", "Dominate",
         [("Total Users", "1,000,000"), ("PRO Users (10%)", "100,000"),
          ("Monthly Recurring Revenue", "$250,000"), ("Affiliate + B2B/mo", "$50,000"),
          ("Annual Revenue", "~$3.6M")], GREEN),
    ]

    for i, (year, phase, metrics, color) in enumerate(years):
        x = Inches(0.4) + i * Inches(4.25)
        box(s, x, Inches(2.1), Inches(4.0), Inches(5.0), fill_color=MID_NAVY)
        box(s, x, Inches(2.1), Inches(4.0), Inches(0.6), fill_color=color)
        txt(s, year, x + Inches(0.15), Inches(2.13), Inches(3.7), Inches(0.4),
            size=Pt(16), bold=True, color=DEEP_NAVY)
        txt(s, phase, x + Inches(0.15), Inches(2.72), Inches(3.7), Inches(0.35),
            size=Pt(12), color=color, italic=True)
        for j, (label, val) in enumerate(metrics):
            y_m = Inches(3.15) + j * Inches(0.75)
            txt(s, label, x + Inches(0.2), y_m, Inches(2.4), Inches(0.4),
                size=Pt(11), color=LIGHT_GREY)
            vc = GOLD if j == len(metrics) - 1 else WHITE
            txt(s, val, x + Inches(0.2), y_m + Inches(0.28), Inches(3.6), Inches(0.4),
                size=Pt(14), bold=(j == len(metrics) - 1), color=vc)

    txt(s, "⚡  Path to profitability: Month 8 | Break-even: 2,500 PRO users",
        Inches(0.6), Inches(7.05), Inches(12), Inches(0.4),
        size=Pt(13), color=GREEN, bold=True)


def slide_gtm(prs):
    s = add_slide(prs)
    bg(s, DEEP_NAVY)
    box(s, 0, 0, W, Inches(0.08), fill_color=ELECTRIC)
    slide_number(s, 9)

    txt(s, "GO-TO-MARKET STRATEGY", Inches(0.6), Inches(0.3), Inches(7), Inches(0.5),
        size=Pt(11), color=ELECTRIC, bold=True)
    txt(s, "Bottom-Up Viral Growth in High-Density SMB Markets",
        Inches(0.6), Inches(0.7), Inches(12), Inches(0.8),
        size=Pt(26), bold=True, color=WHITE)

    phases = [
        ("Phase 1\nMonth 1-3", "India Launch", [
            "Target: Mumbai/Bangalore WhatsApp groups for small traders",
            "Hindi-language onboarding — first tool in Hindi on Telegram",
            "Partner with 5 local micro-influencers (YouTube: small biz tips)",
            "Viral loop: every ad footer says 'Created with AdBot — try free'",
        ], ELECTRIC),
        ("Phase 2\nMonth 4-8", "Africa + LatAm", [
            "Nigeria (Lagos — largest WhatsApp SMB market in Africa)",
            "Brazil (Portuguese — 2nd largest WhatsApp market globally)",
            "Partner: local chambers of commerce + startup hubs",
            "Swahili support unlocks Kenya, Tanzania, Uganda",
        ], GOLD),
        ("Phase 3\nMonth 9-18", "Global + Enterprise", [
            "Arabic launch → Gulf + North Africa (high-spend SMBs)",
            "WhatsApp Business API rollout (enterprise clients)",
            "Agency white-label program ($200/mo per agency seat)",
            "App Store listing: Telegram miniapp for discoverability",
        ], GREEN),
    ]

    for i, (phase, title, items, color) in enumerate(phases):
        x = Inches(0.4) + i * Inches(4.25)
        box(s, x, Inches(1.8), Inches(4.0), Inches(5.3), fill_color=MID_NAVY)
        box(s, x, Inches(1.8), Inches(4.0), Inches(0.7), fill_color=color)
        txt(s, phase, x + Inches(0.15), Inches(1.83), Inches(3.7), Inches(0.55),
            size=Pt(11), bold=True, color=DEEP_NAVY)
        txt(s, title, x + Inches(0.15), Inches(2.55), Inches(3.7), Inches(0.45),
            size=Pt(14), bold=True, color=color)
        for j, item in enumerate(items):
            txt(s, "▸ " + item, x + Inches(0.15), Inches(3.1) + j * Inches(0.75),
                Inches(3.7), Inches(0.7), size=Pt(11), color=WHITE)


def slide_traction(prs):
    s = add_slide(prs)
    bg(s, DEEP_NAVY)
    box(s, 0, 0, W, Inches(0.08), fill_color=GOLD)
    slide_number(s, 10)

    txt(s, "TRACTION — WHAT'S ALREADY BUILT", Inches(0.6), Inches(0.3), Inches(9), Inches(0.5),
        size=Pt(11), color=GOLD, bold=True)
    txt(s, "Working Product. Zero External Capital.",
        Inches(0.6), Inches(0.7), Inches(10), Inches(0.8),
        size=Pt(30), bold=True, color=WHITE)

    built = [
        ("✅ LIVE", "AI Ad Copy Engine", "Gemini 2.5 Flash generates professional copy in 50+ languages with few-shot industry examples"),
        ("✅ LIVE", "Image Generation Pipeline", "Vision DNA: Gemini extracts product DNA from photos → FLUX model generates matched ads"),
        ("✅ LIVE", "Freemium Monetisation", "Telegram Stars payment flow working end-to-end. PRO tier gates Google Ads + Print Poster"),
        ("✅ LIVE", "Rating & Feedback Loop", "Users rate 1-5, bot regenerates with specific feedback. Continuous improvement per user"),
        ("✅ LIVE", "Canva Affiliate Integration", "Every ad delivery includes affiliate-tagged Canva link for passive commission income"),
        ("🔨 Q2", "Background Removal (rembg)", "Real product extraction — replaces Pollinations entirely. Actual product in the ad"),
        ("🔨 Q2", "WhatsApp Business API", "Same bot, new channel. Opens 2.8B user market"),
        ("🔨 Q3", "Pillow Ad Compositor", "Pillow renders headline + CTA on the image — AdCreative.ai quality, $0 cost"),
    ]

    for i, (status, title, desc) in enumerate(built):
        col = 0 if i % 2 == 0 else 1
        row = i // 2
        x = Inches(0.4) + col * Inches(6.5)
        y = Inches(1.8) + row * Inches(1.3)
        color = GREEN if "LIVE" in status else GOLD
        box(s, x, y, Inches(6.2), Inches(1.15), fill_color=MID_NAVY)
        txt(s, status, x + Inches(0.15), y + Inches(0.1), Inches(1.1), Inches(0.4),
            size=Pt(10), bold=True, color=color)
        txt(s, title, x + Inches(1.3), y + Inches(0.1), Inches(4.7), Inches(0.4),
            size=Pt(13), bold=True, color=WHITE)
        txt(s, desc, x + Inches(0.15), y + Inches(0.55), Inches(5.9), Inches(0.55),
            size=Pt(11), color=LIGHT_GREY)


def slide_investment(prs):
    s = add_slide(prs)
    bg(s, DEEP_NAVY)
    box(s, 0, 0, W, Inches(0.08), fill_color=GOLD)
    slide_number(s, 11)

    txt(s, "INVESTMENT ASK", Inches(0.6), Inches(0.3), Inches(5), Inches(0.5),
        size=Pt(11), color=GOLD, bold=True)
    txt(s, "Seed Round: $150,000",
        Inches(0.6), Inches(0.7), Inches(10), Inches(1.0),
        size=Pt(42), bold=True, color=WHITE)
    txt(s, "18-month runway to product-market fit and first $500K ARR",
        Inches(0.6), Inches(1.65), Inches(11), Inches(0.55),
        size=Pt(18), color=LIGHT_GREY)
    accent_line(s, Inches(0.6), Inches(2.3), Inches(6), GOLD)

    # Use of funds
    txt(s, "USE OF FUNDS", Inches(0.6), Inches(2.5), Inches(5), Inches(0.4),
        size=Pt(13), bold=True, color=GOLD)

    funds = [
        ("Engineering & Product", "$70,000", "47%", "Full-time dev 12mo: background removal, compositor, WhatsApp API, multilingual"),
        ("Growth & Marketing", "$40,000", "27%", "Influencer seeding in India, Nigeria, Brazil. Community building in target markets"),
        ("Infrastructure & APIs", "$20,000", "13%", "Server hosting, Gemini API costs at scale, WhatsApp Business API, monitoring"),
        ("Legal & Company Setup", "$12,000", "8%",  "Company incorporation, IP protection, terms of service, privacy policy (GDPR)"),
        ("Working Capital Reserve", "$8,000",  "5%",  "Buffer for unexpected costs and opportunities"),
    ]

    for i, (cat, amt, pct, desc) in enumerate(funds):
        y = Inches(3.0) + i * Inches(0.84)
        box(s, Inches(0.5), y, Inches(12.3), Inches(0.75), fill_color=MID_NAVY)
        # Progress bar
        bar_w = float(pct.strip('%')) / 100 * Inches(5)
        box(s, Inches(0.5), y + Inches(0.62), bar_w, Inches(0.1), fill_color=GOLD)
        txt(s, cat, Inches(0.7), y + Inches(0.05), Inches(2.8), Inches(0.35),
            size=Pt(13), bold=True, color=WHITE)
        txt(s, amt + f"  ({pct})", Inches(3.6), y + Inches(0.05), Inches(1.5), Inches(0.35),
            size=Pt(13), bold=True, color=GOLD)
        txt(s, desc, Inches(5.2), y + Inches(0.05), Inches(7.4), Inches(0.55),
            size=Pt(10), color=LIGHT_GREY)


def slide_roadmap(prs):
    s = add_slide(prs)
    bg(s, DEEP_NAVY)
    box(s, 0, 0, W, Inches(0.08), fill_color=ELECTRIC)
    slide_number(s, 12)

    txt(s, "PRODUCT ROADMAP", Inches(0.6), Inches(0.3), Inches(5), Inches(0.5),
        size=Pt(11), color=ELECTRIC, bold=True)
    txt(s, "3 Phases. 18 Months. Global Scale.",
        Inches(0.6), Inches(0.7), Inches(10), Inches(0.8),
        size=Pt(30), bold=True, color=WHITE)

    phases = [
        ("PHASE 1", "Q2-Q3 2026", "Core Quality", ELECTRIC, [
            "Background removal (rembg)", "Pillow ad compositor (text overlay)",
            "Hindi + Arabic + Spanish launch", "WhatsApp Business API",
            "Template library (6 designs)", "Target: 25,000 users",
        ]),
        ("PHASE 2", "Q4 2026", "Scale", GOLD, [
            "Telegram MiniApp (in-chat UI)", "Brand kit (save colors + logo)",
            "Multi-format batch generation", "Agency white-label tier",
            "Video ad stubs (slideshow)", "Target: 200,000 users",
        ]),
        ("PHASE 3", "Q1-Q2 2027", "Enterprise", GREEN, [
            "fal.ai studio-grade img2img", "Custom AI model fine-tuning",
            "CRM integrations (Shopify, WooCommerce)", "Analytics dashboard",
            "Series A fundraise readiness", "Target: 1,000,000 users",
        ]),
    ]

    for i, (phase, period, title, color, items) in enumerate(phases):
        x = Inches(0.4) + i * Inches(4.25)
        box(s, x, Inches(1.8), Inches(4.0), Inches(5.4), fill_color=MID_NAVY)
        box(s, x, Inches(1.8), Inches(4.0), Inches(0.8), fill_color=color)
        txt(s, phase, x + Inches(0.15), Inches(1.83), Inches(2.5), Inches(0.4),
            size=Pt(14), bold=True, color=DEEP_NAVY)
        txt(s, period, x + Inches(2.5), Inches(1.9), Inches(1.35), Inches(0.35),
            size=Pt(10), color=DEEP_NAVY, align=PP_ALIGN.RIGHT)
        txt(s, title, x + Inches(0.15), Inches(2.65), Inches(3.7), Inches(0.4),
            size=Pt(14), bold=True, color=color)
        for j, item in enumerate(items):
            txt(s, "▸ " + item, x + Inches(0.15), Inches(3.15) + j * Inches(0.68),
                Inches(3.7), Inches(0.63), size=Pt(12), color=WHITE)


def slide_vision(prs):
    s = add_slide(prs)
    bg(s, DEEP_NAVY)
    box(s, 0, 0, W, Inches(0.08), fill_color=GOLD)
    slide_number(s, 13)

    txt(s, "THE VISION", Inches(0.6), Inches(0.3), Inches(5), Inches(0.5),
        size=Pt(11), color=GOLD, bold=True)

    txt(s, '"Every small business owner on Earth',
        Inches(0.6), Inches(1.2), Inches(12), Inches(0.9),
        size=Pt(34), bold=True, color=WHITE)
    txt(s, "can afford a professional marketing team —",
        Inches(0.6), Inches(2.05), Inches(12), Inches(0.9),
        size=Pt(34), bold=False, color=WHITE)
    txt(s, 'in their language, on their phone."',
        Inches(0.6), Inches(2.9), Inches(12), Inches(0.9),
        size=Pt(34), bold=True, color=GOLD)

    accent_line(s, Inches(0.6), Inches(3.95), Inches(11.5), ELECTRIC)

    txt(s, "We are not building another ad tool.\nWe are building the first AI marketing agency that works for the other 95%.",
        Inches(0.6), Inches(4.15), Inches(12), Inches(1.0),
        size=Pt(17), color=LIGHT_GREY)

    items = [
        ("$0", "Cost to small business owner"),
        ("400M+", "Addressable businesses, Day 1"),
        ("50+", "Languages at launch"),
        ("2", "Channels: WhatsApp + Telegram"),
    ]
    for i, (val, label) in enumerate(items):
        label_value(s, label, val, Inches(0.5) + i * Inches(3.2), Inches(5.5),
                    w=Inches(3.0), value_color=ELECTRIC)


def slide_closing(prs):
    s = add_slide(prs)
    bg(s, DEEP_NAVY)
    box(s, 0, 0, W, Inches(0.08), fill_color=ELECTRIC)
    box(s, 0, H - Inches(0.08), W, Inches(0.08), fill_color=GOLD)

    txt(s, "AdBot", Inches(1), Inches(1.5), Inches(11), Inches(1.5),
        size=Pt(72), bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    txt(s, "Let's bring professional advertising to every small business on Earth.",
        Inches(1), Inches(3.0), Inches(11), Inches(0.8),
        size=Pt(20), color=ELECTRIC, align=PP_ALIGN.CENTER)

    accent_line(s, Inches(4), Inches(3.95), Inches(5), GOLD)

    txt(s, "Seed Round: $150,000  |  18-Month Runway  |  Path to $3.6M ARR",
        Inches(1), Inches(4.2), Inches(11), Inches(0.6),
        size=Pt(16), bold=True, color=GOLD, align=PP_ALIGN.CENTER)

    txt(s, "📧  your@email.com     |     📱  Telegram: @YourHandle     |     🌐  adbot.app",
        Inches(1), Inches(5.2), Inches(11), Inches(0.6),
        size=Pt(14), color=LIGHT_GREY, align=PP_ALIGN.CENTER)

    txt(s, "Thank You",
        Inches(1), Inches(5.9), Inches(11), Inches(1.2),
        size=Pt(48), bold=True, color=WHITE, align=PP_ALIGN.CENTER)


# ── Build the deck ────────────────────────────────────────────────────────────

def build():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H

    print("Building slides...")
    slide_cover(prs);        print("  [1/13] Cover ✓")
    slide_problem(prs);      print("  [2/13] Problem ✓")
    slide_solution(prs);     print("  [3/13] Solution ✓")
    slide_market(prs);       print("  [4/13] Market Opportunity ✓")
    slide_pipeline(prs);     print("  [5/13] How It Works ✓")
    slide_competitive(prs);  print("  [6/13] Competitive Landscape ✓")
    slide_business_model(prs); print("  [7/13] Business Model ✓")
    slide_financials(prs);   print("  [8/13] Financial Projections ✓")
    slide_gtm(prs);          print("  [9/13] Go-to-Market ✓")
    slide_traction(prs);     print(" [10/13] Traction ✓")
    slide_investment(prs);   print(" [11/13] Investment Ask ✓")
    slide_roadmap(prs);      print(" [12/13] Roadmap ✓")
    slide_vision(prs);       print(" [13/13] Vision ✓")
    slide_closing(prs);      print(" [14/13] Closing Slide ✓")

    os.makedirs("docs", exist_ok=True)
    out = "docs/AdBot_Investor_Pitch_Deck.pptx"
    prs.save(out)
    print(f"\n✅ Saved: {out}")
    print(f"   Slides: {len(prs.slides)}")
    import os as _os
    size_kb = _os.path.getsize(out) // 1024
    print(f"   Size:   {size_kb} KB")
    print("\n📌 Before presenting:")
    print("   1. Open in PowerPoint / Keynote / Google Slides")
    print("   2. Fill in your contact details on the last slide")
    print("   3. Add your Telegram bot screenshot on slide 3")
    print("   4. Add your real email/LinkedIn on closing slide")


if __name__ == "__main__":
    build()
